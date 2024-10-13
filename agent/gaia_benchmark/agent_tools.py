from langchain_core.tools import tool
from typing import Optional
from agent.gaia_benchmark.gaia_tools.browser import SimpleTextBrowser
import requests
import re
import mimetypes
from markdownify import MarkdownConverter
from youtube_transcript_api import YouTubeTranscriptApi
from urllib.parse import urlparse, parse_qs
import pdfminer.high_level
import mammoth
import pptx
from bs4 import BeautifulSoup
import base64
from PIL import Image
from io import BytesIO
import os
import uuid
import yaml
LLM_ENV = yaml.safe_load(open('config/llm_env.yml', 'r'))

# Initialize the browser instance
browser = SimpleTextBrowser()
if not os.path.isdir("dump/artifacts"):
    os.makedirs("dump/artifacts")

# Helper function to retrieve browser state
def _browser_state() -> str:
    header = f"Address: {browser.address}\n"
    if browser.page_title is not None:
        header += f"Title: {browser.page_title}\n"
    
    current_page = browser.viewport_current_page
    total_pages = len(browser.viewport_pages)
    header += f"Viewport position: Showing page {current_page+1} of {total_pages}."
    return header

@tool()
def web_search(query: str, filter_year: Optional[int] = None) -> str:
    """Perform a web search query and return the search results."""
    browser.visit_page(f"google: {query}", filter_year=filter_year)
    header = _browser_state()
    return header + "\n=======================\n" + browser.viewport
@tool()
def navigational_web_search(query: str) -> str:
    """Perform a navigational web search and visit the top result."""
    browser.visit_page(f"google: {query}")
    
    # Extract the first link and navigate to it
    match = re.search(r"\[.*?\]\((http.*?)\)", browser.page_content)
    if match:
        browser.visit_page(match.group(1))
    
    header = _browser_state()
    return header + "\n=======================\n" + browser.viewport

@tool()
def visit_page(url: str) -> str:
    """Visit a webpage at a given URL and return its text."""  
    browser.visit_page(url)
    header = _browser_state()
    return header + "\n=======================\n" + browser.viewport

@tool()
def download_file(url: str) -> str:
    """Download a file at a given URL (use visit_page for PDF/TXT/HTML files)."""
    if "arxiv" in url:
        url = url.replace("abs", "pdf")
    response = requests.get(url)
    content_type = response.headers.get("content-type", "")
    extension = mimetypes.guess_extension(content_type)
    
    if extension:
        file_path = f"dump/artifacts/file{extension}"
    else:
        file_path = "dump/artifacts/file.object"
    
    with open(file_path, "wb") as f:
        f.write(response.content)

    if "pdf" in extension or "txt" in extension or "html" in extension or 'htm' in extension:
        file_path = browser._fetch_page(url, extract_path=True)
        #raise Exception(f"You tried to use `download_file` with the following:{url}.\n  You must use `visit_page` tool for PDF/TXT/HTML files! Using `download_file` for these file types is not allowed and results an error.")
    
    return f"File was downloaded and saved under path {file_path}. If the file is a PDF/TXT/HTML file, use the `visit_page` tool to view its contents."

@tool()
def page_up() -> str:
    """Scroll the viewport UP one page-length in the current webpage."""
    browser.page_up()
    header = _browser_state()
    return header + "\n=======================\n" + browser.viewport

@tool()
def page_down() -> str:
    """Scroll the viewport DOWN one page-length in the current webpage."""
    browser.page_down()
    header = _browser_state()
    return header + "\n=======================\n" + browser.viewport

@tool()
def find_on_page(search_string: str) -> str:
    """Find the first occurrence of a string on the current page (Ctrl+F)."""
    find_result = browser.find_on_page(search_string)
    header = _browser_state()
    
    if find_result is None:
        return header + f"\n=======================\nThe search string '{search_string}' was not found on this page."
    else:
        return header + "\n=======================\n" + browser.viewport

@tool()
def find_next() -> str:
    """Find the next occurrence of the previously searched string."""
    find_result = browser.find_next()
    header = _browser_state()
    
    if find_result is None:
        return header + "\n=======================\nThe search string was not found on this page."
    else:
        return header + "\n=======================\n" + browser.viewport

# @tool()
# def find_archived_url(url: str, date: str) -> str:
#     """Search Wayback Machine for the archived version of a URL for a given date."""
#     archive_url = f"https://archive.org/wayback/available?url={url}&timestamp={date}"
#     response = requests.get(archive_url).json()
#     try:
#         closest = response["archived_snapshots"]["closest"]
#     except KeyError:
#         raise Exception(f"URL was not archived on Wayback Machine for the given date.")
#
#     browser.visit_page(closest["url"])
#     header = _browser_state()
#     return f"Web archive for URL {url}, snapshot taken at date {closest['timestamp'][:8]}:\n" + header + "\n=======================\n" + browser.viewport

# tools from mdconvert.py

@tool()
def plain_text_conversion(file_path: str, file_extension: Optional[str] = None) -> str:
    """Convert a plain text file into a text string."""
    if file_extension == "":
        return "Error: Invalid file extension."

    content_type, encoding = mimetypes.guess_type(f"__placeholder{file_extension}")

    with open(file_path, "rt") as file_handle:
        text_content = file_handle.read()

    return text_content

@tool()
def html_conversion(file_path: str, file_extension: Optional[str] = None) -> str:
    """Convert an HTML file into markdown text."""
    if file_extension.lower() not in [".html", ".htm"]:
        return "Error: Invalid file extension for HTML."

    with open(file_path, "rt") as file_handle:
        html_content = file_handle.read()

    soup = BeautifulSoup(html_content, "html.parser")
    
    # Remove script and style blocks
    for script in soup(["script", "style"]):
        script.extract()

    body_elm = soup.find("body")
    markdown_text = ""
    if body_elm:
        markdown_text = MarkdownConverter().convert_soup(body_elm)
    else:
        markdown_text = MarkdownConverter().convert_soup(soup)

    return markdown_text

@tool()
def youtube_transcript(url: str) -> str:
    """Extract the transcript of a YouTube video using the video URL."""
    parsed_url = urlparse(url)
    params = parse_qs(parsed_url.query)

    if "v" not in params:
        return "Error: Invalid YouTube URL."

    video_id = params["v"][0]
    transcript = YouTubeTranscriptApi.get_transcript(video_id)

    transcript_text = " ".join([part["text"] for part in transcript])

    return transcript_text


@tool()
def pdf_conversion(file_path: str, file_extension: Optional[str] = None) -> str:
    """Convert a PDF file into a text string."""
    text_content = pdfminer.high_level.extract_text(file_path)
    return text_content

@tool()
def docx_conversion(file_path: str, file_extension: Optional[str] = None) -> str:
    """Convert a DOCX file into a markdown text."""
    if file_extension.lower() != ".docx":
        return "Error: Invalid file extension for DOCX."

    with open(file_path, "rb") as docx_file:
        result = mammoth.convert_to_html(docx_file)
        html_content = result.value
        soup = BeautifulSoup(html_content, "html.parser")
        markdown_text = MarkdownConverter().convert_soup(soup)

    return markdown_text

@tool()
def pptx_conversion(file_path: str, file_extension: Optional[str] = None) -> str:
    """Convert a PPTX file into markdown text."""
    if file_extension.lower() != ".pptx":
        return "Error: Invalid file extension for PPTX."

    md_content = ""
    presentation = pptx.Presentation(file_path)
    slide_num = 0
    for slide in presentation.slides:
        slide_num += 1
        md_content += f"\n\n<!-- Slide number: {slide_num} -->\n"

        title = slide.shapes.title
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape == title:
                    md_content += "# " + shape.text.lstrip() + " "
                else:
                    md_content += shape.text + " "

    return md_content

def encode_image(image_path: str) -> str:
    with open(image_path, "rb") as image_file:
        encoded_string = base64.b64encode(image_file.read()).decode("utf-8")
    return encoded_string


@tool()
def visualizer(image_path: str, question: Optional[str] = None) -> str:
    """A tool that can answer questions about attached images."""
    
    add_note = False
    if not question:
        add_note = True
        question = "Please write a detailed caption for this image."

    base64_image = encode_image(image_path)

    payload = {
        "model": "gpt-4o-mini",
        "messages": [
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": question
                    },
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:image/jpeg;base64,{base64_image}"
                        }
                    }
                ]
            }
        ],
        "max_tokens": 500
    }
    ENDPOINT = f"{LLM_ENV['azure']['AZURE_OPENAI_ENDPOINT']}/openai/deployments/gpt-4o-mini/chat/completions?api-version=2024-02-15-preview"
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {LLM_ENV['azure']['AZURE_OPENAI_API_KEY']}"
    }
    response = requests.post(ENDPOINT, headers=headers, json=payload)
    try:
        output = response.json()['choices'][0]['message']['content']
    except Exception:
        raise Exception(f"Response format unexpected: {response.json()}")

    if add_note:
        output = f"You did not provide a particular question, so here is a detailed caption for the image: {output}"

    return output

@tool
def parse_yaml_code(yaml_code: str) -> dict:
    """You must use this tool before sending the final output, the input is the yaml code with the output schema. The result is the final output!"""
    return "The Yaml doesn't have a valid yaml structure, please fix it such that it can be parsed. Remember that if you have a value that is a string, you should wrap it in quotes."
